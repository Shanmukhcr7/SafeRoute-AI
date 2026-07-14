from pptx import Presentation

prs = Presentation("AI_Road_Accident_RiskZone_Detection_Review1_UPDATED_FIXED.pptx")
for i, slide in enumerate(prs.slides):
    print(f"--- Slide {i+1} ---")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        print(shape.text.encode('utf-8', 'replace').decode('utf-8'))
