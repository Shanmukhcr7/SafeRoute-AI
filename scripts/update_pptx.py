from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

def replace_text_in_pptx(filename, output_filename, replacements):
    prs = Presentation(filename)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            # Read through paragraphs and runs
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for old_text, new_text in replacements.items():
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
                            shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                            shape.text_frame.word_wrap = True
                            
    prs.save(output_filename)
    print("Fixed PPT saved as", output_filename)

if __name__ == "__main__":
    replacements = {
        "govt. open data / Kaggle RTA records": "France Accidents 2005-2016 Kaggle Dataset",
        "Govt / Kaggle RTA dataset": "France Accidents 2005-2016 (Kaggle)"
    }
    replace_text_in_pptx("AI_Road_Accident_RiskZone_Detection_Review1_UPDATED_FIXED.pptx", "AI_Road_Accident_RiskZone_Detection_Review1_FINAL.pptx", replacements)
