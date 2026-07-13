from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

def fix_pptx(filename, output_filename):
    prs = Presentation(filename)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            text_frame = shape.text_frame
            # Enable auto-fit text
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            text_frame.word_wrap = True
            
            # Optional: reduce font size slightly for runs to ensure they fit
            # But auto_size should ideally handle it if the PPT viewer supports it
            
    prs.save(output_filename)
    print("Fixed PPT saved as", output_filename)

if __name__ == "__main__":
    fix_pptx("AI_Road_Accident_RiskZone_Detection_Review1_UPDATED.pptx", "AI_Road_Accident_RiskZone_Detection_Review1_UPDATED_FIXED.pptx")
